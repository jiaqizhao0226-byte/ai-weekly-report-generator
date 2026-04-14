#!/usr/bin/env python3
"""
AI News PPT Generator - Fixed version without complex slide duplication
"""

from flask import Flask, render_template, request, jsonify, send_file
from flask_cors import CORS
from datetime import datetime
import os
import sys

# Load .env file
from dotenv import load_dotenv
load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))

sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from news_fetcher import format_news_for_display, get_ai_news, expand_news_batch, rewrite_news_batch, search_wechat_articles

app = Flask(__name__)
CORS(app)

def generate_news_text_parts(news_item):
    """Split news into title (bold) and content (normal) parts."""
    rewritten = news_item.get('rewritten', '')
    if rewritten:
        # 清理换行符，统一为一段话
        rewritten = rewritten.replace('\n', ' ').replace('\r', ' ').strip()

        if '：' in rewritten:
            parts = rewritten.split('：', 1)
            title_part = parts[0].strip() + '：'
            content_part = parts[1].strip() if len(parts) > 1 else ''
        elif ':' in rewritten:
            parts = rewritten.split(':', 1)
            title_part = parts[0].strip() + '：'
            content_part = parts[1].strip() if len(parts) > 1 else ''
        else:
            # 没有冒号，取前25字当标题
            title_part = rewritten[:25].strip() + '：'
            content_part = rewritten[25:].strip()
    else:
        title = news_item.get('title', '')
        description = news_item.get('description', '')
        date = news_item.get('date_full', news_item.get('date', ''))

        if date and '-' in date and '年' not in date:
            try:
                dt = datetime.strptime(date[:10], '%Y-%m-%d')
                date = dt.strftime('%Y年%m月%d日')
            except:
                pass

        title_part = title + '：'
        content_part = f"{date}，{description}" if date else description

    # 确保title_part不为空
    if not title_part or title_part == '：':
        title_part = content_part[:25] + '：'
        content_part = content_part[25:]

    # Soft length limit - cut at last complete sentence if too long
    max_total = 320
    total = len(title_part) + len(content_part)
    if total > max_total:
        # Find the last sentence-ending punctuation within limit
        max_content = max_total - len(title_part)
        truncated = content_part[:max_content]
        # Cut at last 。or ；
        for punct in ['。', '；', '，']:
            last_pos = truncated.rfind(punct)
            if last_pos > len(truncated) // 2:  # At least keep half
                content_part = truncated[:last_pos + 1]
                break
        else:
            content_part = truncated

    return title_part, content_part

def download_image(url, timeout=10):
    """Download image from URL and return as BytesIO"""
    try:
        import requests
        from io import BytesIO
        from urllib.parse import urlparse

        parsed = urlparse(url)
        referer = f"{parsed.scheme}://{parsed.netloc}/"

        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
            'Referer': referer,
            'Accept': 'image/webp,image/apng,image/*,*/*;q=0.8'
        }
        response = requests.get(url, headers=headers, timeout=timeout)
        if response.status_code == 200 and len(response.content) > 1000:
            return BytesIO(response.content)
    except:
        pass
    return None

def search_logo_image(news_item, timeout=10):
    """Search for company/model logo when article has no image"""
    import requests
    from io import BytesIO

    title = news_item.get('title', '')

    logo_keywords = [
        'OpenAI', 'Anthropic', 'Google', 'Microsoft', 'Meta', 'Apple', 'Amazon', 'Nvidia',
        '百度', '阿里', '腾讯', '字节跳动', '华为', '小米', '商汤', '科大讯飞',
        'ChatGPT', 'GPT', 'Claude', 'Gemini', 'Llama', 'Mistral', 'Qwen', '千问',
        'DeepSeek', 'Kimi', '文心', '混元', '豆包', 'Sora', 'Midjourney', 'Stable Diffusion'
    ]

    search_term = None
    for kw in logo_keywords:
        if kw.lower() in title.lower():
            search_term = f"{kw} logo"
            break

    if not search_term:
        search_term = title.split('：')[0] if '：' in title else title[:20]
        search_term = f"{search_term} logo AI"

    brave_api_key = os.environ.get('BRAVE_API_KEY', '')
    if not brave_api_key:
        return None

    try:
        url = "https://api.search.brave.com/res/v1/images/search"
        headers = {
            "Accept": "application/json",
            "X-Subscription-Token": brave_api_key
        }
        params = {
            "q": search_term,
            "count": 5,
            "safesearch": "strict"
        }

        response = requests.get(url, headers=headers, params=params, timeout=timeout)
        if response.status_code == 200:
            data = response.json()
            results = data.get('results', [])

            for result in results:
                img_url = result.get('properties', {}).get('url') or result.get('thumbnail', {}).get('src')
                if img_url:
                    img_stream = download_image(img_url, timeout=5)
                    if img_stream:
                        print(f"  Found logo for '{search_term}': {img_url[:60]}...")
                        return img_stream
    except Exception as e:
        print(f"  Logo search error: {e}")

    return None

def replace_shape_with_image(slide, shape, image_stream):
    """Replace a text shape with an image at the same position, preserving aspect ratio"""
    try:
        from PIL import Image
        from io import BytesIO

        left = shape.left
        top = shape.top
        placeholder_width = shape.width
        placeholder_height = shape.height

        image_stream.seek(0)
        img = Image.open(image_stream)
        img_width, img_height = img.size

        if img.format not in ['BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF']:
            png_stream = BytesIO()
            if img.mode in ('RGBA', 'P'):
                img = img.convert('RGBA')
                img.save(png_stream, format='PNG')
            else:
                img = img.convert('RGB')
                img.save(png_stream, format='PNG')
            png_stream.seek(0)
            image_stream = png_stream
        else:
            image_stream.seek(0)

        img_ratio = img_width / img_height
        placeholder_ratio = placeholder_width / placeholder_height

        if img_ratio > placeholder_ratio:
            new_width = placeholder_width
            new_height = int(placeholder_width / img_ratio)
        else:
            new_height = placeholder_height
            new_width = int(placeholder_height * img_ratio)

        new_left = left + (placeholder_width - new_width) // 2
        new_top = top + (placeholder_height - new_height) // 2

        sp = shape._element
        sp.getparent().remove(sp)

        slide.shapes.add_picture(image_stream, new_left, new_top, new_width, new_height)
        return True
    except Exception as e:
        print(f"  Image insert error: {e}")
        return False

def delete_slide(prs, index):
    """Delete slide at index"""
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]

def fill_slide_news(slide, news_list):
    """Fill a slide with news items + images. Supports both 1-item and 2-item slides."""
    news_shapes = []
    image_shapes = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if '新闻' in text:
                news_shapes.append(shape)
            elif '图片占位' in text:
                image_shapes.append(shape)

    # Sort top to bottom
    news_shapes.sort(key=lambda s: s.top)
    image_shapes.sort(key=lambda s: s.top)

    for i, shape in enumerate(news_shapes):
        if i < len(news_list):
            news_item = news_list[i]
            title_part, content_part = generate_news_text_parts(news_item)
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]

            # 统一行距
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.3

            # Title (bold)
            run1 = p.add_run()
            run1.text = title_part
            run1.font.size = Pt(14)
            run1.font.bold = True
            run1.font.name = '华文楷体'
            run1.font.color.rgb = RGBColor(0, 0, 0)

            # Content
            run2 = p.add_run()
            run2.text = content_part
            run2.font.size = Pt(14)
            run2.font.bold = False
            run2.font.name = '华文楷体'
            run2.font.color.rgb = RGBColor(0, 0, 0)

            # Fill corresponding image placeholder
            if i < len(image_shapes):
                image_url = news_item.get('image', '')
                img_stream = None
                if image_url:
                    img_stream = download_image(image_url)
                if not img_stream:
                    img_stream = search_logo_image(news_item)
                if img_stream:
                    replace_shape_with_image(slide, image_shapes[i], img_stream)
                else:
                    image_shapes[i].text_frame.clear()
                    p = image_shapes[i].text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "[暂无图片]"
                    run.font.size = Pt(10)
                    run.font.name = '微软雅黑'
                    run.font.color.rgb = RGBColor(128, 128, 128)
        else:
            # No news for this slot, clear text and image
            shape.text_frame.clear()
            if i < len(image_shapes):
                image_shapes[i].text_frame.clear()

def generate_ppt(selected_news, output_path, date_range):
    """Generate PPT: fill news into 14-slide template, delete unused slides.

    Template layout:
    0: cover
    1-3: model double (2 news each) + 4: model single (1 news)
    5-7: application double + 8: application single
    9-11: investment double + 12: investment single
    13: end
    """
    template_path = os.path.join(os.path.dirname(__file__), 'AI周报模板.pptx')
    prs = Presentation(template_path)

    # Update cover date (slide 0)
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if '2025年' in text or '光子' in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = date_range + '  |  光子策略分析'
                run.font.size = Pt(14)
                run.font.name = '华文楷体'
                run.font.color.rgb = RGBColor(255, 255, 255)

    # Group by category
    news_by_cat = {'model': [], 'application': [], 'investment': []}
    for news in selected_news:
        cat = news.get('category', 'application')
        if cat in news_by_cat:
            news_by_cat[cat].append(news)

    cat_config = {
        'model':       {'double': [1, 2, 3], 'single': 4},
        'application': {'double': [5, 6, 7], 'single': 8},
        'investment':  {'double': [9, 10, 11], 'single': 12},
    }

    slides_to_delete = []

    for cat in ['model', 'application', 'investment']:
        config = cat_config[cat]
        double_slides = config['double']
        single_slide = config['single']
        cat_news = news_by_cat[cat]

        n = len(cat_news)
        pairs = n // 2
        has_odd = n % 2 == 1

        # Fill double-slides
        for i, slide_idx in enumerate(double_slides):
            if i < pairs:
                news_start = i * 2
                fill_slide_news(prs.slides[slide_idx], cat_news[news_start:news_start + 2])
            else:
                slides_to_delete.append(slide_idx)

        # Fill or delete single-slide
        if has_odd:
            fill_slide_news(prs.slides[single_slide], [cat_news[-1]])
        else:
            slides_to_delete.append(single_slide)

    # Delete unused slides in reverse order
    for idx in sorted(slides_to_delete, reverse=True):
        delete_slide(prs, idx)

    prs.save(output_path)
    return output_path

# Flask routes
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/search', methods=['POST'])
def api_search():
    data = request.json
    days = data.get('days', 7)
    news = get_ai_news(days)
    news = format_news_for_display(news)
    news.sort(key=lambda x: x.get('importance', 0), reverse=True)
    return jsonify({'news': news, 'count': len(news)})

@app.route('/api/keyword-search', methods=['POST'])
def api_keyword_search():
    """关键词搜索公众号文章"""
    data = request.json
    keywords = data.get('keywords', '')
    days = data.get('days', 7)

    if not keywords or not keywords.strip():
        return jsonify({'success': False, 'error': '请输入搜索关键词'}), 400

    try:
        results = search_wechat_articles(keywords, days=days)
        results = format_news_for_display(results)
        return jsonify({'success': True, 'news': results, 'count': len(results)})
    except Exception as e:
        import traceback
        return jsonify({'success': False, 'error': str(e), 'trace': traceback.format_exc()}), 500

@app.route('/api/fetch-url', methods=['POST'])
def api_fetch_url():
    """Fetch article from URL for manual addition"""
    import requests
    from bs4 import BeautifulSoup
    import re

    data = request.json
    url = data.get('url', '')
    category = data.get('category', 'application')

    if not url:
        return jsonify({'success': False, 'error': '请提供链接'})

    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        response = requests.get(url, headers=headers, timeout=15)
        response.encoding = response.apparent_encoding
        soup = BeautifulSoup(response.text, 'html.parser')

        title = ''
        for selector in ['h1', 'title', '.article-title', '.post-title']:
            elem = soup.select_one(selector)
            if elem:
                title = elem.get_text().strip()
                if len(title) > 10:
                    break

        description = ''
        meta_desc = soup.select_one('meta[name="description"]')
        if meta_desc:
            description = meta_desc.get('content', '')

        if not description or len(description) < 50:
            for selector in ['article p', '.article-content p', '.post-content p', 'main p']:
                paragraphs = soup.select(selector)
                if paragraphs:
                    texts = [p.get_text().strip() for p in paragraphs[:3] if len(p.get_text().strip()) > 30]
                    if texts:
                        description = ' '.join(texts)[:500]
                        break

        date_str = ''
        for selector in ['.date', '.publish-time', '.post-date', 'time[datetime]', '.article-time']:
            elem = soup.select_one(selector)
            if elem:
                text = elem.get('datetime', '') or elem.text.strip()
                date_match = re.search(r'(\d{4})[年/-](\d{1,2})[月/-](\d{1,2})', text)
                if date_match:
                    date_str = f"{date_match.group(1)}-{date_match.group(2).zfill(2)}-{date_match.group(3).zfill(2)}"
                    break

        image = ''
        og_img = soup.select_one('meta[property="og:image"]')
        if og_img and og_img.get('content'):
            image = og_img.get('content')

        from urllib.parse import urlparse
        source = urlparse(url).netloc

        if date_str:
            try:
                dt = datetime.strptime(date_str, '%Y-%m-%d')
                date_display = dt.strftime('%m月%d日')
                date_full = dt.strftime('%Y年%m月%d日')
            except:
                date_display = '近期'
                date_full = ''
        else:
            now = datetime.now()
            date_display = now.strftime('%m月%d日')
            date_full = now.strftime('%Y年%m月%d日')
            date_str = now.strftime('%Y-%m-%d')

        news_item = {
            'title': title or '未知标题',
            'description': description,
            'url': url,
            'source': source,
            'date': date_str,
            'date_display': date_display,
            'date_full': date_full,
            'image': image,
            'category': category,
            'importance': 75,
            'custom': True
        }

        return jsonify({'success': True, 'news': news_item})

    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/api/generate', methods=['POST'])
def api_generate():
    data = request.json
    selected_news = data.get('news', [])
    date_range = data.get('date_range', datetime.now().strftime('%Y年%m月第%W周'))

    # AI rewrite news in professional tone (max 220 chars each)
    selected_news = rewrite_news_batch(selected_news, max_chars=180)

    output_dir = os.path.join(os.path.dirname(__file__), 'output')
    os.makedirs(output_dir, exist_ok=True)

    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    output_path = os.path.join(output_dir, f'AI_Weekly_Report_{timestamp}.pptx')

    try:
        generate_ppt(selected_news, output_path, date_range)
        return jsonify({'success': True, 'filename': os.path.basename(output_path)})
    except Exception as e:
        import traceback
        return jsonify({'success': False, 'error': str(e), 'trace': traceback.format_exc()}), 500

@app.route('/api/download')
def api_download():
    filename = request.args.get('file', '')
    output_dir = os.path.join(os.path.dirname(__file__), 'output')

    if filename:
        file_path = os.path.join(output_dir, filename)
    else:
        files = sorted([f for f in os.listdir(output_dir) if f.endswith('.pptx')], reverse=True)
        file_path = os.path.join(output_dir, files[0]) if files else None

    if file_path and os.path.exists(file_path):
        return send_file(file_path, as_attachment=True, download_name=os.path.basename(file_path))
    return jsonify({'error': 'File not found'}), 404

if __name__ == '__main__':
    os.makedirs(os.path.join(os.path.dirname(__file__), 'templates'), exist_ok=True)
    os.makedirs(os.path.join(os.path.dirname(__file__), 'output'), exist_ok=True)
    print("Starting AI News PPT Generator on http://localhost:5050")
    app.run(host='0.0.0.0', port=5050, debug=True)
